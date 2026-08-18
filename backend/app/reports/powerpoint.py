from __future__ import annotations
from io import BytesIO
from .schemas import Report

class PowerPointReportGenerator:
    def generate(self, report: Report) -> bytes:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("PowerPoint generation requires the 'python-pptx' package.") from exc
        presentation = Presentation(); slide = presentation.slides.add_slide(presentation.slide_layouts[0]); slide.shapes.title.text = report.title; slide.placeholders[1].text = report.executive_summary
        for heading, text in (("KPIs", "\n".join(f"{item['name']}: {item['value']}" for item in report.kpis)), ("Insights & Recommendations", "\n".join([item.get("description", "") for item in report.insights] + report.recommendations)), ("Forecast", str(report.forecast) if report.forecast else "No validated forecast available.")):
            page = presentation.slides.add_slide(presentation.slide_layouts[1]); page.shapes.title.text = heading; page.placeholders[1].text = text or "No data available."
        buffer = BytesIO(); presentation.save(buffer); return buffer.getvalue()
